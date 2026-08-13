from pathlib import Path
from pptx.enum.shapes import MSO_SHAPE_TYPE
from io import BytesIO
from PIL import Image
from pptx import Presentation
from Documents.model import Document
from Documents.ocr import extract_text, ask_ocr
from Material.config import CONSOLE


def extract_slide_text(slide) -> str:
    """ extracts all selectable text from a PowerPoint slide """

    text = []

    for shape in slide.shapes:
        if not shape.has_text_frame: continue

        for paragraph in shape.text_frame.paragraphs:

            content = paragraph.text.strip()
            if content: text.append(content)

    return "\n".join(text)


def extract_slide_images(slide, width, height, min_area: float = 0.05) -> list[dict]:
    """ extracts embedded images from a PowerPoint slide """

    images = []
    slide_area = width * height

    for shape in slide.shapes:
        
        if shape.shape_type != MSO_SHAPE_TYPE.PICTURE: continue

        area_ratio = (shape.width * shape.height) / slide_area
        if area_ratio < min_area: continue

        image = shape.image
        images.append({"bytes": image.blob, "extension": image.ext})

    return images


def import_pptx(path: Path, dpi: int, language: str) -> Document | None:
    """Imports a PowerPoint presentation and optionally uses OCR."""

    try: presentation = Presentation(path)
    except Exception as e: return CONSOLE.print(f"[red]import_documents: Error {e}[/red]")

    width = presentation.slide_width
    height = presentation.slide_height

    text = []; slide_information = []
    has_text = False; has_images = False; ocr_used = False; use_ocr = False

    slides = len(presentation.slides)

    # goes through first to see needs and ocr information
    for slide in presentation.slides:

        slide_text = extract_slide_text(slide)
        images = extract_slide_images(slide, width, height)

        if slide_text: has_text = True
        if images: has_images = True

        slide_information.append((slide_text, images))

    has_ocr_candidates = any(images and not slide_text for slide_text, images in slide_information)

    if has_ocr_candidates: use_ocr = ask_ocr(path, slides)

    # build the final document content.
    for slide_number, (slide_text, images) in enumerate(slide_information, start = 1):
        slide_content = []

        if slide_text: slide_content.append(slide_text)     # native PowerPoint text

        if images and use_ocr:      # OCR embedded images when necessary

            for image in images:
                # Convert image bytes into something pytesseract/PIL can process.
                image_object = Image.open(BytesIO(image["bytes"]))

                ocr_text = extract_text(image_object, language)

                if ocr_text:
                    slide_content.append(ocr_text)
                    ocr_used = True

        if slide_content: text.append(f"Slide {slide_number}\n" + "\n".join(slide_content))

    # Create the common Document object.
    title = (path.stem.replace("_", " ").replace("-", " "))
    return Document(title, "\n\n".join(text), slides, has_text, has_images, ocr_used)